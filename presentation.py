from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import os
from presentationContent import presentationContent


def create_presentation(data):
    presentation = Presentation()

    for slide_data in data:
        slide_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(slide_layout)

        image_path = slide_data['Image']
        image = Image.open(image_path)
        image_width, image_height = image.size

        if image_width > image_height:
            max_width = Inches(5)
            image_ratio = image_height / image_width
            height = max_width * image_ratio
            width = max_width
        else:
            max_height = Inches(6)
            image_ratio = image_width / image_height
            width = max_height * image_ratio
            height = max_height

        left = Inches(0.5)
        top = (presentation.slide_height - height) / 2  # Center the image vertically
        converted_image_path = convert_to_jpeg(image_path)
        slide.shapes.add_picture(converted_image_path, left, top, width, height)

        text_box_width = Inches(10) - width - Inches(1)  # Calculate the width of the text box
        text_box_left = left + width + Inches(0.25)  # Calculate the left position of the text box

        text_box = slide.shapes.add_textbox(left=text_box_left, top=0, width=text_box_width, height=height)
        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.word_wrap = True

        header = text_frame.add_paragraph()
        header.text = slide_data['Header']
        header.font.size = Pt(24)
        header.font.bold = True

        content = text_frame.add_paragraph()
        content.text = slide_data['Content']
        content.font.size = Pt(18)

    return presentation


def convert_to_jpeg(image_path):
    image = Image.open(image_path)
    rgb_image = image.convert("RGB")  # Convert to RGB mode
    converted_image_path = f'{image_path}.jpg'
    rgb_image.save(converted_image_path, 'JPEG')
    return converted_image_path


# Example usage
query = "Elon Musk"
pages = 10
max_num = 40
data = presentationContent(query, pages, max_num)
presentation = create_presentation(data)
presentation.save(f"{query}.pptx")

for file in os.listdir("Images"):
    os.remove(f"Images/{file}")

